"""
Movie Recommendation System - PPT Generator
Generates a professional, template-based PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── COLOR PALETTE ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)   # Deep navy  – slide background
ACCENT_RED   = RGBColor(0xE5, 0x09, 0x14)   # Official Netflix red – headings / accents
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xD0, 0xD0, 0xD0)   # Body text
MID_GRAY     = RGBColor(0x88, 0x88, 0x88)   # Sub-labels / captions
CARD_BG      = RGBColor(0x16, 0x21, 0x3E)   # Slightly lighter navy for cards
ACCENT_LINE  = RGBColor(0xE5, 0x00, 0x00)   # Thin separator line

# ─── SLIDE DIMENSIONS (16:9 widescreen) ───────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── HELPERS ──────────────────────────────────────────────────────────────────

def set_bg(slide, color=DARK_BG):
    """Fill slide background with a solid colour."""
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()   # no border
    return shape

def add_accent_line(slide, left, top, width, height=Pt(2)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_RED
    shape.line.fill.background()
    return shape

def add_bullet_textbox(slide, bullets, left, top, width, height,
                       font_size=16, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.space_before = Pt(4)
    return txBox

def slide_header(slide, title, subtitle=None):
    """Standard top-bar header used on every content slide."""
    # Red left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT_RED)
    # Title
    add_text(slide, title,
             left=Inches(0.3), top=Inches(0.18),
             width=Inches(12.5), height=Inches(0.65),
             font_size=30, bold=True, color=WHITE)
    # Red underline
    add_accent_line(slide, Inches(0.3), Inches(0.88), Inches(12.5))
    if subtitle:
        add_text(slide, subtitle,
                 left=Inches(0.3), top=Inches(0.95),
                 width=Inches(12.5), height=Inches(0.4),
                 font_size=15, color=MID_GRAY)

def add_card(slide, left, top, width, height, title, body_lines,
             title_size=17, body_size=14):
    """A rounded-corner-look card (rectangle + text)."""
    add_rect(slide, left, top, width, height, CARD_BG)
    # Card title
    add_text(slide, title,
             left=left + Inches(0.15), top=top + Inches(0.1),
             width=width - Inches(0.3), height=Inches(0.4),
             font_size=title_size, bold=True, color=ACCENT_RED)
    # Thin red separator inside card
    add_accent_line(slide, left + Inches(0.15), top + Inches(0.52),
                    width - Inches(0.3), Pt(1.5))
    # Card body
    add_bullet_textbox(slide, body_lines,
                       left=left + Inches(0.15),
                       top=top + Inches(0.62),
                       width=width - Inches(0.3),
                       height=height - Inches(0.75),
                       font_size=body_size, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 – Title / Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    # Big red left stripe
    add_rect(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT_RED)

    # Main title
    add_text(slide, "Movie Recommendation System",
             left=Inches(0.5), top=Inches(1.8),
             width=Inches(12.5), height=Inches(1.2),
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Sub-title
    add_text(slide, "A Content-Based Filtering Approach Using TF-IDF & Multiple ML Models",
             left=Inches(0.5), top=Inches(3.1),
             width=Inches(11.5), height=Inches(0.85),
             font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    add_accent_line(slide, Inches(0.5), Inches(3.0), Inches(5))

    # Dataset info pill
    add_rect(slide, Inches(0.5), Inches(4.1), Inches(3.5), Inches(0.55), CARD_BG)
    add_text(slide, "📦  Dataset : movies_metadata.csv  |  45,000+ movies",
             left=Inches(0.6), top=Inches(4.15),
             width=Inches(5), height=Inches(0.45),
             font_size=13, color=MID_GRAY)

    # Decorative bottom strip
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), ACCENT_RED)


def slide_overview(prs):
    """Slide 2 – Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Project Overview",
                 "What are we building and why?")

    # 3 info cards side by side
    card_data = [
        ("🎯  Problem Statement",
         ["Users are overwhelmed by thousands",
          "of movie choices.",
          "",
          "We need a smart system that",
          "recommends similar movies based",
          "on what a user already likes."]),
        ("💡  Our Solution",
         ["Content-Based Filtering using",
          "TF-IDF vectorization on movie",
          "metadata (overview + genres +",
          "tagline).",
          "",
          "Four ML models are trained &",
          "compared for accuracy."]),
        ("📊  Dataset",
         ["Source: TMDB movies_metadata.csv",
          "Size : 45,466 movies (raw)",
          "After cleaning : ~45,447 movies",
          "",
          "Features used:",
          "• title   • overview",
          "• genres  • tagline"]),
    ]
    cw = Inches(3.9)
    ch = Inches(4.7)
    for i, (t, b) in enumerate(card_data):
        add_card(slide,
                 left=Inches(0.3 + i * 4.35), top=Inches(1.55),
                 width=cw, height=ch,
                 title=t, body_lines=b, title_size=16, body_size=13)


def slide_pipeline(prs):
    """Slide 3 – Implementation Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Implementation Pipeline",
                 "Step-by-step workflow from raw data to recommendations")

    steps = [
        ("1", "Load Data", "Read CSV\nmovies_metadata.csv"),
        ("2", "Clean Data", "Drop duplicates\nHandle nulls"),
        ("3", "Feature Eng.", "Combine overview\n+ genres + tagline → tags"),
        ("4", "NLP Pre-proc.", "Lowercase, remove punct.,\nstopwords, lemmatize"),
        ("5", "TF-IDF", "Vectorize tags\n(50k features, 1-2 ngrams)"),
        ("6", "4 ML Models", "Cosine · KNN\nSigmoid · K-Means"),
        ("7", "Evaluate", "Confusion Matrix\n& Metrics"),
    ]

    box_w = Inches(1.6)
    box_h = Inches(2.5)
    gap   = Inches(0.18)
    start_x = Inches(0.22)
    top_y = Inches(1.75)

    for i, (num, title, desc) in enumerate(steps):
        lx = start_x + i * (box_w + gap)
        # Card background
        add_rect(slide, lx, top_y, box_w, box_h, CARD_BG)
        # Number bubble (red)
        add_rect(slide, lx + Inches(0.55), top_y + Inches(0.1),
                 Inches(0.5), Inches(0.45), ACCENT_RED)
        add_text(slide, num,
                 left=lx + Inches(0.55), top=top_y + Inches(0.1),
                 width=Inches(0.5), height=Inches(0.45),
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_text(slide, title,
                 left=lx + Inches(0.08), top=top_y + Inches(0.65),
                 width=box_w - Inches(0.16), height=Inches(0.45),
                 font_size=13, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
        # Desc
        add_text(slide, desc,
                 left=lx + Inches(0.08), top=top_y + Inches(1.15),
                 width=box_w - Inches(0.16), height=Inches(1.2),
                 font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        # Arrow between boxes
        if i < len(steps) - 1:
            add_text(slide, "→",
                     left=lx + box_w, top=top_y + Inches(1.0),
                     width=gap, height=Inches(0.5),
                     font_size=18, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)

    # Bottom note
    add_text(slide,
             "★  TF-IDF Vector Model:  Tags are transformed into sparse 45447×50000 matrix – the foundation for ALL four models below.",
             left=Inches(0.3), top=Inches(4.5),
             width=Inches(12.5), height=Inches(0.6),
             font_size=13, italic=True, color=MID_GRAY)


def slide_tfidf(prs):
    """Slide 4 – TF-IDF Vector Model Deep Dive"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Vector Representation : TF-IDF",
                 "The core text model powering all four recommendation algorithms")

    left_bullets = [
        "• TF-IDF = Term Frequency × Inverse Document Frequency",
        "",
        "• Assigns higher weight to rare but important words",
        "  (e.g., 'dystopian', 'heist') and downweights common",
        "  words (e.g., 'movie', 'story').",
        "",
        "• Input  → combined 'tags' column  (overview + genres + tagline)",
        "• Output → sparse matrix  [45,447 movies × 50,000 features]",
        "",
        "• Settings used:",
        "   max_features = 50,000",
        "   ngram_range  = (1, 2)   ← unigrams + bigrams",
        "   stop_words   = 'english'",
    ]

    right_bullets = [
        "Why TF-IDF over simple Bag-of-Words?",
        "",
        "✔  Eliminates noisy, high-frequency words",
        "✔  Surfaces genre-defining, plot-specific terms",
        "✔  Produces compact, meaningful numeric vectors",
        "✔  Scales efficiently to 45k+ documents",
        "",
        "Result Matrix Shape:",
        "  (45447,  50000)",
        "",
        "Stored as:   tfidf_matrix.pkl",
        "Vectorizer:  tfidf.pkl",
    ]

    # Left panel
    add_rect(slide, Inches(0.25), Inches(1.55), Inches(6.3), Inches(5.4), CARD_BG)
    add_text(slide, "How it works",
             left=Inches(0.4), top=Inches(1.6),
             width=Inches(6), height=Inches(0.4),
             font_size=16, bold=True, color=ACCENT_RED)
    add_accent_line(slide, Inches(0.4), Inches(2.05), Inches(6.0))
    add_bullet_textbox(slide, left_bullets,
                       left=Inches(0.4), top=Inches(2.15),
                       width=Inches(6.0), height=Inches(4.6),
                       font_size=13, color=LIGHT_GRAY)

    # Right panel
    add_rect(slide, Inches(6.8), Inches(1.55), Inches(6.2), Inches(5.4), CARD_BG)
    add_text(slide, "Advantages & Stats",
             left=Inches(6.95), top=Inches(1.6),
             width=Inches(5.9), height=Inches(0.4),
             font_size=16, bold=True, color=ACCENT_RED)
    add_accent_line(slide, Inches(6.95), Inches(2.05), Inches(5.9))
    add_bullet_textbox(slide, right_bullets,
                       left=Inches(6.95), top=Inches(2.15),
                       width=Inches(5.9), height=Inches(4.6),
                       font_size=13, color=LIGHT_GRAY)


def slide_models_summary(prs):
    """Slide 5 – All Four Models at a Glance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Models Used  —  Snapshot",
                 "Four distinct algorithms built on the same TF-IDF representation")

    models = [
        ("Model 1", "Cosine Similarity",
         "Computes angle between\ntwo TF-IDF vectors.\nRange: 0 (different) → 1 (identical).\nFastest & most common."),
        ("Model 2", "KNN\n(K-Nearest Neighbors)",
         "Finds 10 nearest neighbours\nin vector space using cosine\ndistance.\nScikit-learn BruteForce, n_jobs=-1."),
        ("Model 3", "Sigmoid Kernel\n(SVM-based)",
         "Non-linear kernel similarity\nfrom kernel trick.\nsigmoid_kernel() applied\nrow-vs-matrix."),
        ("Model 4", "K-Means Clustering",
         "Groups all 45k movies into\n20 clusters.\nRecommends movies from the\nsame cluster as input."),
    ]

    cw = Inches(2.95)
    ch = Inches(4.5)
    for i, (tag, name, desc) in enumerate(models):
        lx = Inches(0.3 + i * 3.22)
        add_rect(slide, lx, Inches(1.7), cw, ch, CARD_BG)
        # Tag bullet
        add_rect(slide, lx, Inches(1.7), cw, Inches(0.35), ACCENT_RED)
        add_text(slide, tag,
                 left=lx, top=Inches(1.7),
                 width=cw, height=Inches(0.35),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, name,
                 left=lx + Inches(0.1), top=Inches(2.1),
                 width=cw - Inches(0.2), height=Inches(0.75),
                 font_size=15, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
        add_accent_line(slide, lx + Inches(0.1), Inches(2.88), cw - Inches(0.2))
        add_text(slide, desc,
                 left=lx + Inches(0.1), top=Inches(2.98),
                 width=cw - Inches(0.2), height=Inches(3.0),
                 font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)


def slide_model_detail(prs, num, name, how_it_works, formula_lines, result_sample):
    """Template for per-model detail slides (Slides 6–9)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, f"Model {num}: {name}")

    # Left card – How it works
    add_rect(slide, Inches(0.25), Inches(1.35), Inches(5.5), Inches(5.7), CARD_BG)
    add_text(slide, "🔧  How it works",
             left=Inches(0.4), top=Inches(1.4),
             width=Inches(5.2), height=Inches(0.4),
             font_size=16, bold=True, color=ACCENT_RED)
    add_accent_line(slide, Inches(0.4), Inches(1.85), Inches(5.2))
    add_bullet_textbox(slide, how_it_works,
                       left=Inches(0.4), top=Inches(1.95),
                       width=Inches(5.2), height=Inches(4.9),
                       font_size=13, color=LIGHT_GRAY)

    # Middle card – Formula / Key concept
    add_rect(slide, Inches(6.0), Inches(1.35), Inches(3.5), Inches(5.7), CARD_BG)
    add_text(slide, "📐  Key Formula",
             left=Inches(6.15), top=Inches(1.4),
             width=Inches(3.2), height=Inches(0.4),
             font_size=16, bold=True, color=ACCENT_RED)
    add_accent_line(slide, Inches(6.15), Inches(1.85), Inches(3.2))
    add_bullet_textbox(slide, formula_lines,
                       left=Inches(6.15), top=Inches(1.95),
                       width=Inches(3.2), height=Inches(4.9),
                       font_size=13, color=LIGHT_GRAY)

    # Right card – Sample output for "Avatar"
    add_rect(slide, Inches(9.7), Inches(1.35), Inches(3.35), Inches(5.7), CARD_BG)
    add_text(slide, "📋  Sample Output  (Avatar)",
             left=Inches(9.85), top=Inches(1.4),
             width=Inches(3.1), height=Inches(0.55),
             font_size=14, bold=True, color=ACCENT_RED)
    add_accent_line(slide, Inches(9.85), Inches(1.97), Inches(3.1))
    add_bullet_textbox(slide, result_sample,
                       left=Inches(9.85), top=Inches(2.07),
                       width=Inches(3.1), height=Inches(4.8),
                       font_size=12, color=LIGHT_GRAY)


def slide_cosine(prs):
    how = [
        "1. Compute cosine similarity between the",
        "   query movie vector and ALL other movie",
        "   vectors in the TF-IDF matrix.",
        "",
        "2. Sort similarity scores in descending order.",
        "",
        "3. Return top-10 most similar movies",
        "   (excluding the query itself).",
        "",
        "• Library  : sklearn.metrics.pairwise",
        "             cosine_similarity()",
        "• Time complexity : O(n × d)",
        "  n = movies, d = TF-IDF features",
    ]
    formula = [
        "cos(θ) = (A · B) / (‖A‖ × ‖B‖)",
        "",
        "• A, B  →  TF-IDF row vectors",
        "• ·     →  dot product",
        "• ‖·‖   →  L2 norm",
        "",
        "Score range: 0 → 1",
        "  0 = completely different",
        "  1 = identical",
        "",
        "Threshold: none required —",
        "rank by score directly.",
    ]
    results = [
        "1. Avatar 2",
        "2. The Inhabited Island",
        "3. Thor: Ragnarok",
        "4. Moontrap: Target Earth",
        "5. The Three Musketeers",
        "6. A Trip to the Moon",
        "7. Nightmare City 2035",
        "8. France société anonyme",
        "9. Désiré",
        "10. Stand by Me Doraemon",
    ]
    slide_model_detail(prs, "1", "Cosine Similarity  (TF-IDF)", how, formula, results)


def slide_knn(prs):
    how = [
        "1. Fit NearestNeighbors model on the",
        "   TF-IDF matrix (all 45k movies).",
        "",
        "2. For the query movie, find k+1",
        "   nearest neighbours using cosine",
        "   distance (brute-force).",
        "",
        "3. Exclude the query itself and",
        "   return the top-10 titles.",
        "",
        "• sklearn.neighbors.NearestNeighbors",
        "• n_neighbors=11  (10 recs + self)",
        "• metric='cosine'",
        "• algorithm='brute'",
        "• n_jobs=-1  (all CPU cores)",
    ]
    formula = [
        "distance = 1 − cosine_similarity",
        "",
        "k-NN retrieves k movies with",
        "smallest cosine distance from",
        "the query vector.",
        "",
        "Why brute-force?",
        "  Sparse high-dim vectors don't",
        "  benefit from tree-based index.",
        "",
        "Output: indices + distances",
        "sorted ascending by distance.",
    ]
    results = [
        "1. Avatar 2",
        "2. The Inhabited Island",
        "3. Thor: Ragnarok",
        "4. Moontrap: Target Earth",
        "5. The Three Musketeers",
        "6. A Trip to the Moon",
        "7. Désiré",
        "8. France société anonyme",
        "9. Nightmare City 2035",
        "10. Stand by Me Doraemon",
    ]
    slide_model_detail(prs, "2", "K-Nearest Neighbors (KNN)", how, formula, results)


def slide_sigmoid(prs):
    how = [
        "1. Compute the sigmoid kernel between",
        "   the query movie vector and ALL",
        "   movie vectors.",
        "",
        "2. Sort sigmoid scores descending.",
        "",
        "3. Return the top-10 results.",
        "",
        "• Library  : sklearn.metrics.pairwise",
        "             sigmoid_kernel()",
        "• Inspired by SVM kernel trick —",
        "  captures non-linear similarity",
        "  relationships between documents.",
    ]
    formula = [
        "K(x, y) = tanh(γ·xᵀy + c)",
        "",
        "• tanh  →  hyperbolic tangent",
        "• γ     →  kernel coefficient",
        "           (default=1/n_features)",
        "• c     →  intercept (default=1)",
        "",
        "Produces a non-linear",
        "similarity score — can surface",
        "subtler content relationships",
        "missed by cosine similarity.",
    ]
    results = [
        "1. Avatar 2",
        "2. The Inhabited Island",
        "3. Thor: Ragnarok",
        "4. Moontrap: Target Earth",
        "5. The Three Musketeers",
        "6. A Trip to the Moon",
        "7. Nightmare City 2035",
        "8. France société anonyme",
        "9. Désiré",
        "10. Stand by Me Doraemon",
    ]
    slide_model_detail(prs, "3", "Sigmoid Kernel Similarity", how, formula, results)


def slide_kmeans(prs):
    how = [
        "1. Apply K-Means clustering to the",
        "   TF-IDF matrix to form 20 clusters.",
        "",
        "2. Every movie is assigned a cluster_id",
        "   (numeric label 0-19).",
        "",
        "3. For the query movie, retrieve its",
        "   cluster_id, then return the first",
        "   10 movies in that same cluster.",
        "",
        "• sklearn.cluster.KMeans",
        "• n_clusters = 20",
        "• random_state = 42",
        "• No distance ranking — just",
        "  cluster membership.",
    ]
    formula = [
        "Minimize  Σ ‖xᵢ − μⱼ‖²",
        "",
        "• xᵢ  →  TF-IDF vector of movie i",
        "• μⱼ  →  centroid of cluster j",
        "",
        "Each doc assigned to nearest",
        "centroid; centroids updated",
        "iteratively until convergence.",
        "",
        "20 clusters ≈ broad genre /",
        "theme groups across 45k movies.",
    ]
    results = [
        "1. Powder",
        "2. The City of Lost Children",
        "3. Twelve Monkeys",
        "4. Lawnmower Man 2",
        "5. Screamers",
        "6. Shopping",
        "7. Congo",
        "8. Johnny Mnemonic",
        "9. Judge Dredd",
        "10. Power Rangers: The Movie",
    ]
    slide_model_detail(prs, "4", "K-Means Clustering", how, formula, results)


def slide_evaluation_method(prs):
    """Slide – Evaluation Methodology"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Evaluation Methodology",
                 "How we measure 'relevance' — Multi-Genre Overlap approach")

    left_bullets = [
        "Test Movies Used:",
        "  1. Avatar          4. Titanic",
        "  2. The Dark Knight  5. Inception",
        "  3. Toy Story",
        "",
        "For each test movie:",
        " • Get top-10 recommendations from each model.",
        " • Sample 10 random non-recommended movies.",
        " • Evaluate all 20 movies as a set.",
        "",
        "Prediction label (y_pred):",
        "  1  if the movie was recommended",
        "  0  if it was NOT recommended",
    ]

    right_bullets = [
        "Ground Truth (y_true) — Genre Overlap:",
        "",
        "• Extract genres of the input movie",
        "  as a Python set.",
        "  e.g., Avatar → {'Action', 'Adventure',",
        "                   'Science Fiction'}",
        "",
        "• For each candidate movie extract",
        "  its genres as a set.",
        "",
        "• y_true = 1  if  intersection ≠ ∅",
        "  (at least one shared genre)",
        "",
        "• y_true = 0  otherwise",
        "",
        "✔ Fairer than primary-genre-only match",
    ]

    add_rect(slide, Inches(0.25), Inches(1.55), Inches(6.1), Inches(5.5), CARD_BG)
    add_text(slide, "Test Setup",
             left=Inches(0.4), top=Inches(1.6),
             width=Inches(5.8), height=Inches(0.4),
             font_size=16, bold=True, color=ACCENT_RED)
    add_accent_line(slide, Inches(0.4), Inches(2.05), Inches(5.8))
    add_bullet_textbox(slide, left_bullets,
                       left=Inches(0.4), top=Inches(2.15),
                       width=Inches(5.8), height=Inches(4.75),
                       font_size=13, color=LIGHT_GRAY)

    add_rect(slide, Inches(6.6), Inches(1.55), Inches(6.45), Inches(5.5), CARD_BG)
    add_text(slide, "Relevance Ground Truth",
             left=Inches(6.75), top=Inches(1.6),
             width=Inches(6.2), height=Inches(0.4),
             font_size=16, bold=True, color=ACCENT_RED)
    add_accent_line(slide, Inches(6.75), Inches(2.05), Inches(6.2))
    add_bullet_textbox(slide, right_bullets,
                       left=Inches(6.75), top=Inches(2.15),
                       width=Inches(6.2), height=Inches(4.75),
                       font_size=13, color=LIGHT_GRAY)


def slide_confusion_matrices(prs, img_path):
    """Slide – Confusion Matrices (with embedded image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Confusion Matrices  —  All 4 Models",
                 "Evaluated over 5 test movies  ×  20 candidate movies each")

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 Inches(0.3), Inches(1.35),
                                 Inches(12.6), Inches(5.8))
    else:
        add_text(slide, f"[Image not found: {img_path}]",
                 left=Inches(1), top=Inches(3),
                 width=Inches(11), height=Inches(1),
                 font_size=16, color=ACCENT_RED, align=PP_ALIGN.CENTER)


def slide_final_results(prs):
    """Slide – Final Metrics Table & Winner"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Final Results  —  Model Comparison",
                 "Accuracy · Precision · Recall · F1-Score across all 4 models")

    # ── Table header ──────────────────────────────────────────
    col_headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score"]
    col_x       = [Inches(0.3), Inches(4.0), Inches(6.2), Inches(8.4), Inches(10.6)]
    col_w       = [Inches(3.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.5)]
    row_h       = Inches(0.52)
    header_y    = Inches(1.55)

    # Header bar
    add_rect(slide, Inches(0.3), header_y, Inches(12.5), row_h, ACCENT_RED)
    for h, lx, w in zip(col_headers, col_x, col_w):
        add_text(slide, h,
                 left=lx, top=header_y,
                 width=w, height=row_h,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Data rows ──────────────────────────────────────────────
    rows = [
        ("Cosine Similarity",   "70.0 %", "90.0 %", "64.29 %", "75.00 %"),
        ("KNN",                 "70.0 %", "90.0 %", "64.29 %", "75.00 %"),
        ("Sigmoid Kernel",      "70.0 %", "90.0 %", "64.29 %", "75.00 %"),
        ("K-Means Clustering",  "72.0 %", "90.0 %", "66.18 %", "76.27 %"),  # winner
    ]

    winner_idx = 3   # K-Means (0-indexed)
    for r_idx, row in enumerate(rows):
        row_y = header_y + row_h + r_idx * row_h
        bg = RGBColor(0x22, 0x30, 0x52) if r_idx % 2 == 0 else CARD_BG
        if r_idx == winner_idx:
            bg = RGBColor(0x5C, 0x00, 0x00)   # darker red highlight
        add_rect(slide, Inches(0.3), row_y, Inches(12.5), row_h, bg)
        for val, lx, w in zip(row, col_x, col_w):
            fc = WHITE if r_idx != winner_idx else RGBColor(0xFF, 0xDD, 0xDD)
            add_text(slide, val,
                     left=lx, top=row_y,
                     width=w, height=row_h,
                     font_size=13, bold=(r_idx == winner_idx), color=fc,
                     align=PP_ALIGN.CENTER)

    # ── Winner badge ──────────────────────────────────────────
    badge_y = Inches(4.9)
    add_rect(slide, Inches(0.3), badge_y, Inches(12.5), Inches(0.65), CARD_BG)
    add_text(slide,
             "🏆  Best Model by F1-Score:  K-Means Clustering  (F1 = 76.27%)",
             left=Inches(0.5), top=badge_y,
             width=Inches(12.1), height=Inches(0.65),
             font_size=18, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)

    # ── Insight bullets ──────────────────────────────────────
    insights = [
        "★  Cosine, KNN and Sigmoid produce identical results — all are distance-based on the same TF-IDF matrix.",
        "★  K-Means slightly outperforms because cluster-level grouping captures broader thematic similarity.",
        "★  Precision is uniformly high (90 %) — models rarely recommend truly irrelevant movies.",
        "★  Recall gap (64 % vs 66 %) is small; all models miss some genre-matching movies not in top-10.",
    ]
    add_bullet_textbox(slide, insights,
                       left=Inches(0.4), top=Inches(5.65),
                       width=Inches(12.4), height=Inches(1.7),
                       font_size=12, color=MID_GRAY)


def slide_model_comparison_chart(prs, img_path):
    """Slide – Model Comparison Bar Chart (embedded image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_header(slide, "Model Comparison Chart",
                 "F1-Score & Accuracy visualised across all four models")

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 Inches(1.5), Inches(1.35),
                                 Inches(10.0), Inches(5.8))
    else:
        add_text(slide, f"[Image not found: {img_path}]",
                 left=Inches(1), top=Inches(3),
                 width=Inches(11), height=Inches(1),
                 font_size=16, color=ACCENT_RED, align=PP_ALIGN.CENTER)


def slide_conclusion(prs):
    """Last Slide – Conclusion & Key Takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Red left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT_RED)

    add_text(slide, "Conclusion & Key Takeaways",
             left=Inches(0.5), top=Inches(0.3),
             width=Inches(12.5), height=Inches(0.8),
             font_size=34, bold=True, color=WHITE)
    add_accent_line(slide, Inches(0.5), Inches(1.15), Inches(5))

    takeaways = [
        ("📌  TF-IDF + Content-Based Filtering works well",
         "Combining overview, genres and tagline into a unified 'tags' field\nand vectorizing with TF-IDF creates a rich, informative representation."),
        ("📌  All distance-based models perform equally",
         "Cosine, KNN and Sigmoid share the same TF-IDF space and produce\nidentical top-10 results — they differ only in computation style."),
        ("📌  K-Means adds cluster-level diversity",
         "Slightly higher F1 (76.27 %) because it considers global structure\nover the entire corpus, not just pairwise distances."),
        ("📌  Genre Overlap is a better evaluation criterion",
         "Multi-genre intersection (vs. primary-genre exact-match) gives\na fairer, more robust picture of recommendation quality."),
    ]

    cw = Inches(5.9)
    ch = Inches(2.3)
    for i, (title, body) in enumerate(takeaways):
        row = i // 2
        col = i % 2
        lx = Inches(0.4 + col * 6.45)
        ly = Inches(1.4 + row * 2.55)
        add_rect(slide, lx, ly, cw, ch, CARD_BG)
        add_text(slide, title,
                 left=lx + Inches(0.12), top=ly + Inches(0.1),
                 width=cw - Inches(0.24), height=Inches(0.45),
                 font_size=14, bold=True, color=ACCENT_RED)
        add_accent_line(slide, lx + Inches(0.12), ly + Inches(0.58),
                        cw - Inches(0.24))
        add_text(slide, body,
                 left=lx + Inches(0.12), top=ly + Inches(0.68),
                 width=cw - Inches(0.24), height=ch - Inches(0.8),
                 font_size=12, color=LIGHT_GRAY)

    # Bottom strip
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), ACCENT_RED)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    base = os.path.dirname(os.path.abspath(__file__))
    cm_img   = os.path.join(base, "confusion_matrices_all_models.png")
    chart_img = os.path.join(base, "model_comparison_chart.png")

    print("Building slides...")
    slide_title(prs)                        # 1
    slide_overview(prs)                     # 2
    slide_pipeline(prs)                     # 3
    slide_tfidf(prs)                        # 4
    slide_models_summary(prs)               # 5
    slide_cosine(prs)                       # 6
    slide_knn(prs)                          # 7
    slide_sigmoid(prs)                      # 8
    slide_kmeans(prs)                       # 9
    slide_evaluation_method(prs)            # 10
    slide_confusion_matrices(prs, cm_img)   # 11
    slide_model_comparison_chart(prs, chart_img)  # 12
    slide_final_results(prs)                # 13
    slide_conclusion(prs)                   # 14

    out = os.path.join(base, "Movie_Recommendation_System.pptx")
    prs.save(out)
    print(f"\n✅  PPT saved  →  {out}")
    print(f"    Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
